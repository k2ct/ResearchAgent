"""
Document Ingestion Pipeline v1 with optional MinerU backend.

Converts raw documents (PDF, DOCX, PPTX, MD) from raw_docs/ into
YAML-front-matter markdown files in data/ingested/, ready for RAG indexing.

Backends:
- local (default): text extraction with pymupdf / python-docx / python-pptx
- mineru (optional): MinerU API for complex PDFs, tables, formulas, OCR
  Falls back to local on failure.

v1 scope (local backend):
- Text extraction only (no OCR, no image understanding, no formula parsing)
- No LLM calls
"""

import os
from pathlib import Path
from typing import Dict, List, Optional

import yaml

from research_agent.rag.schemas import DIR_TO_SOURCE_TYPE
from research_agent.rag.loaders import extract_title

SUPPORTED_SUFFIXES = {".pdf", ".docx", ".pptx", ".md"}


# ── Text Extractors ──────────────────────────────────────────────


def extract_text_from_pdf(path: Path) -> str:
    """
    Extract text from a PDF file using pymupdf (fitz).

    Each page is prefixed with ``## Page N`` marker.
    """
    try:
        import fitz
    except ImportError:
        raise ImportError(
            "pymupdf is required for PDF extraction. "
            "Install it with: pip install pymupdf"
        )

    try:
        doc = fitz.open(str(path))
    except Exception as e:
        raise ValueError(f"Failed to open PDF: {path}\n{type(e).__name__}: {e}")

    parts: List[str] = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text")
        if text.strip():
            parts.append(f"## Page {i}\n\n{text.strip()}")

    doc.close()

    if not parts:
        return f"## Page 1\n\n(PDF contains no extractable text: {path.name})"

    return "\n\n".join(parts)


def extract_text_from_docx(path: Path) -> str:
    """
    Extract text from a DOCX file using python-docx.

    Paragraphs are joined with double newlines. Empty paragraphs are skipped.
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx is required for DOCX extraction. "
            "Install it with: pip install python-docx"
        )

    try:
        doc = Document(str(path))
    except Exception as e:
        raise ValueError(f"Failed to open DOCX: {path}\n{type(e).__name__}: {e}")

    paragraphs: List[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    if not paragraphs:
        return f"(DOCX contains no extractable text: {path.name})"

    return "\n\n".join(paragraphs)


def extract_text_from_pptx(path: Path) -> str:
    """
    Extract text from a PPTX file using python-pptx.

    Each slide is prefixed with ``## Slide N`` marker.
    Slides with no text content are skipped.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX extraction. "
            "Install it with: pip install python-pptx"
        )

    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise ValueError(f"Failed to open PPTX: {path}\n{type(e).__name__}: {e}")

    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_text = shape.text_frame.text.strip()
                if shape_text:
                    texts.append(shape_text)

        if texts:
            parts.append(f"## Slide {i}\n\n" + "\n\n".join(texts))

    if not parts:
        return f"## Slide 1\n\n(PPTX contains no extractable text: {path.name})"

    return "\n\n".join(parts)


def extract_text_from_md(path: Path) -> str:
    """
    Read a Markdown file directly.
    """
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="gbk", errors="replace")


# ── Metadata Helpers ─────────────────────────────────────────────


def infer_source_type(path: Path) -> str:
    """
    Infer source_type from the parent directory name.

    Uses DIR_TO_SOURCE_TYPE from schemas.py.
    """
    parent_name = path.parent.name
    source_type = DIR_TO_SOURCE_TYPE.get(parent_name)

    if source_type is None:
        return "misc_doc"

    return source_type


def build_metadata(
    path: Path,
    text: str,
    *,
    ingestion_backend: str = "local",
    mineru_used: bool = False,
    mineru_error: str = "",
    mineru_raw_output_path: str = "",
) -> Dict:
    """
    Build a metadata dict for an ingested document.

    Returns a dict suitable for YAML front matter serialization.
    """
    source_type = infer_source_type(path)
    title = extract_title(text) or path.stem

    meta = {
        "source_type": source_type,
        "title": title,
        "doc_type": path.suffix.lstrip("."),
        "original_path": str(path).replace("\\", "/"),
        "created_from": "ingestion_pipeline",
        "ingestion_backend": ingestion_backend,
        "topic": "unknown",
        "tags": [path.parent.name],
    }

    if mineru_used:
        meta["mineru_used"] = True
    if mineru_error:
        meta["mineru_error"] = mineru_error
    if mineru_raw_output_path:
        meta["mineru_raw_output_path"] = mineru_raw_output_path

    return meta


# ── Output Writer ────────────────────────────────────────────────


def write_ingested_markdown(text: str, metadata: Dict, output_path: Path) -> None:
    """
    Write an ingested markdown file with YAML front matter.

    Output format::

        ---
        source_type: paper_note
        title: My Paper
        ...
        ---

        {text}
    """
    output_path.parent.mkdir(parents=True, exist_ok=True)

    yaml_block = yaml.dump(
        metadata,
        allow_unicode=True,
        sort_keys=False,
        default_flow_style=False,
    ).strip()

    full_content = f"---\n{yaml_block}\n---\n\n{text}"

    output_path.write_text(full_content, encoding="utf-8")


# ── Core Ingestion Functions ─────────────────────────────────────


def _resolve_backend(explicit_backend: Optional[str] = None) -> str:
    """
    Resolve the ingestion backend.

    Priority:
    1. Explicit ``backend`` parameter passed to ingest_file / ingest_directory
    2. ``DOCUMENT_INGESTION_BACKEND`` environment variable
    3. Default: ``"local"``
    """
    if explicit_backend:
        return explicit_backend.strip().lower()
    return os.getenv("DOCUMENT_INGESTION_BACKEND", "local").strip().lower()


def _extract_text_local(path: Path, suffix: str) -> str:
    """
    Dispatch to the correct local text extractor based on file suffix.
    """
    if suffix == ".pdf":
        return extract_text_from_pdf(path)
    elif suffix == ".docx":
        return extract_text_from_docx(path)
    elif suffix == ".pptx":
        return extract_text_from_pptx(path)
    elif suffix == ".md":
        return extract_text_from_md(path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def ingest_file(
    path: Path,
    output_dir: Path,
    *,
    backend: Optional[str] = None,
) -> Dict:
    """
    Ingest a single file into the RAG pipeline.

    Args:
        path: Path to the source file (.pdf, .docx, .pptx, .md).
        output_dir: Directory to write the ingested .md file.
        backend: ``"local"`` or ``"mineru"``.  Overrides the
            ``DOCUMENT_INGESTION_BACKEND`` env var.  Defaults to ``"local"``.

    Returns:
        A result dict with keys: status, input_path, output_path,
        source_type, doc_type, title, char_count, error,
        backend, mineru_used, mineru_error, fallback.
    """
    resolved_backend = _resolve_backend(backend)

    result = {
        "status": "error",
        "input_path": str(path),
        "output_path": "",
        "source_type": "",
        "doc_type": path.suffix.lstrip("."),
        "title": "",
        "char_count": 0,
        "error": "",
        "backend": resolved_backend,
        "mineru_used": False,
        "mineru_error": "",
        "fallback": False,
    }

    try:
        suffix = path.suffix.lower()

        # ── MinerU path ──
        if resolved_backend == "mineru":
            from .mineru_client import parse_document_with_mineru

            mineru_result = parse_document_with_mineru(path)

            if mineru_result["ok"]:
                text = mineru_result["markdown_text"]
                metadata = build_metadata(
                    path, text,
                    ingestion_backend="mineru",
                    mineru_used=True,
                    mineru_error="",
                    mineru_raw_output_path=mineru_result.get("raw_output_path", ""),
                )
                result["mineru_used"] = True
            else:
                # Fallback to local extractor
                text = _extract_text_local(path, suffix)
                metadata = build_metadata(
                    path, text,
                    ingestion_backend="local",
                    mineru_used=False,
                    mineru_error=mineru_result.get("mineru_error", "mineru_failed"),
                )
                result["fallback"] = True
                result["mineru_error"] = mineru_result.get("mineru_error", "")
        else:
            text = _extract_text_local(path, suffix)
            metadata = build_metadata(path, text, ingestion_backend="local")

        output_path = output_dir / f"{path.stem}.md"
        write_ingested_markdown(text, metadata, output_path)

        result["status"] = "success"
        result["output_path"] = str(output_path)
        result["source_type"] = metadata["source_type"]
        result["title"] = metadata["title"]
        result["char_count"] = len(text)

    except Exception as e:
        result["error"] = f"{type(e).__name__}: {e}"

    return result


def ingest_directory(
    input_dir: Path,
    output_dir: Path,
    *,
    backend: Optional[str] = None,
) -> List[Dict]:
    """
    Recursively scan input_dir for supported files and ingest them all.

    Skips README.md (case-insensitive). A single file failure does not
    interrupt the batch.

    Returns a list of result dicts, one per file processed.
    """
    if not input_dir.exists():
        raise FileNotFoundError(f"Input directory does not exist: {input_dir}")

    results: List[Dict] = []

    for path in sorted(input_dir.rglob("*")):
        if not path.is_file():
            continue

        if path.name.lower() == "readme.md":
            continue

        if path.suffix.lower() not in SUPPORTED_SUFFIXES:
            continue

        result = ingest_file(path, output_dir, backend=backend)
        results.append(result)

    return results


# ── Reporting Helpers ────────────────────────────────────────────


def print_ingest_summary(results: List[Dict]) -> None:
    """
    Print a human-readable summary of ingestion results.
    """
    if not results:
        print("No files found to ingest.")
        return

    succeeded = [r for r in results if r["status"] == "success"]
    failed = [r for r in results if r["status"] != "success"]
    fallback_count = sum(1 for r in succeeded if r.get("fallback"))
    mineru_count = sum(1 for r in succeeded if r.get("mineru_used"))

    # Detect backend from first result
    backend = results[0].get("backend", "local") if results else "local"

    print(f"\nScanned: {len(results)} files")
    print(f"  Backend: {backend}")
    print(f"  Succeeded: {len(succeeded)}")
    if backend == "mineru":
        print(f"    - via MinerU: {mineru_count}")
        print(f"    - fallback to local: {fallback_count}")
    print(f"  Failed:    {len(failed)}")

    if succeeded:
        print(f"\n--- Succeeded ---")
        for r in succeeded:
            tags = []
            if r.get("mineru_used"):
                tags.append("mineru")
            if r.get("fallback"):
                tags.append("fallback->local")
            tag_str = f" [{', '.join(tags)}]" if tags else ""
            print(f"  {r['source_type']:20s} {r['input_path']}{tag_str}")
            print(f"  {'':20s} -> {r['output_path']} ({r['char_count']} chars)")

    if failed:
        print(f"\n--- Failed ---")
        for r in failed:
            print(f"  {r['input_path']}")
            print(f"  Error: {r['error']}")
